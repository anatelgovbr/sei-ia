"""Regression tests for PPTX files with image shapes without embedded media."""

from __future__ import annotations

from io import BytesIO
from zipfile import ZIP_DEFLATED, ZipFile

import pytest

pytest.importorskip("docling")
pytest.importorskip("pptx")
pytest.importorskip("PIL")

from PIL import Image
from pptx import Presentation

from sei_extraction.parsers.office import extract_with_docling


DRAWINGML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
RELATIONSHIP_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _make_pptx_with_picture(path, *, remove_embed: bool) -> None:
    image = BytesIO()
    Image.new("RGB", (1, 1), color="white").save(image, format="PNG")

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(image, left=0, top=0, width=914400, height=914400)
    text_box = slide.shapes.add_textbox(0, 914400, 9144000, 914400)
    text_box.text_frame.text = "PPTX de reprodução"
    presentation.save(path)

    if not remove_embed:
        return

    with ZipFile(path) as source:
        files = {name: source.read(name) for name in source.namelist()}

    slide_xml = files["ppt/slides/slide1.xml"]
    from lxml import etree

    root = etree.fromstring(slide_xml)
    for blip in root.findall(f".//{{{DRAWINGML_NS}}}blip"):
        blip.attrib.pop(f"{{{RELATIONSHIP_NS}}}embed", None)

    files["ppt/slides/slide1.xml"] = etree.tostring(
        root, xml_declaration=True, encoding="UTF-8", standalone=True
    )
    with ZipFile(path, "w", compression=ZIP_DEFLATED) as target:
        for name, content in files.items():
            target.writestr(name, content)


def test_pptx_without_embedded_image_is_extracted(tmp_path):
    path = tmp_path / "sem-imagem-incorporada.pptx"
    _make_pptx_with_picture(path, remove_embed=True)

    result = extract_with_docling(str(path))

    assert "PPTX de reprodução" in result


def test_pptx_with_embedded_image_is_extracted(tmp_path):
    path = tmp_path / "com-imagem-incorporada.pptx"
    _make_pptx_with_picture(path, remove_embed=False)

    result = extract_with_docling(str(path))

    assert "PPTX de reprodução" in result
